import logging
from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt
from pptx.dml.color import RGBColor
import traceback
from pptx.util import Inches
import re
import math

# 配置日志记录
logging.basicConfig(level=logging.INFO)

class pptTextFormat:
    """存储文本格式信息"""
    def __init__(self, font_name=None, font_size=None, font_bold=None, 
                 font_italic=None, font_underline=None, color=None, 
                 alignment=None, spacing=None, margin_left=None, margin_right=None):
        self.font_name = font_name
        self.font_size = font_size  # 单位pt
        self.font_bold = font_bold
        self.font_italic = font_italic
        self.font_underline = font_underline
        self.color = color  # 新增颜色属性
        self.alignment = alignment
        self.spacing = spacing
        self.margin_left = margin_left
        self.margin_right = margin_right

def split_text(text, max_length=1000):
    """分割长文本"""
    if len(text) <= max_length:
        return [text]
    
    sentences = re.split('([。！？.!?])', text)
    chunks = []
    current_chunk = ""
    for i in range(0, len(sentences), 2):
        sentence = sentences[i]
        if i + 1 < len(sentences):
            sentence += sentences[i + 1]
        if len(current_chunk) + len(sentence) <= max_length:
            current_chunk += sentence
        else:
            if current_chunk:
                chunks.append(current_chunk)
            current_chunk = sentence
    if current_chunk:
        chunks.append(current_chunk)
    return chunks

def calculate_font_size(shape, text, original_size):
    """计算字体大小（保持原逻辑）"""
    try:
        if not shape.width or not shape.height:
            return original_size
        width_pt = shape.width / 12700.0
        height_pt = shape.height / 12700.0
        avg_char_width = original_size * 1.2
        if avg_char_width == 0:
            return original_size
        chars_per_line = width_pt / avg_char_width
        if chars_per_line <= 0:
            return original_size
        estimated_lines = len(text) / chars_per_line
        line_height = original_size * 1.2
        required_height = estimated_lines * line_height
        if required_height > height_pt:
            new_size = original_size * (height_pt / required_height)
            return max(8, new_size)
        return original_size
    except Exception as e:
        logging.warning(f"字体大小计算错误: {e}")
        return original_size

def get_text_format(run, shape=None):
    """获取文本格式信息（增强颜色处理）"""
    try:
        font = run.font
        color_rgb = None
        if font.color and font.color.type == 1:  # RGB颜色类型
            color_rgb = RGBColor(font.color.rgb[0], font.color.rgb[1], font.color.rgb[2])
        elif font.color and hasattr(font.color, 'theme_color'):  # 主题颜色处理
            color_rgb = font.color.theme_color

        format_info = pptTextFormat(
            font_name=font.name,
            font_size=font.size.pt if font.size else None,
            font_bold=font.bold,
            font_italic=font.italic,
            font_underline=font.underline,
            color=color_rgb  # 直接记录颜色对象
        )
        if hasattr(run, '_p'):
            paragraph = run._p
            if hasattr(paragraph, 'alignment'):
                format_info.alignment = paragraph.alignment
            if hasattr(paragraph, 'spacing'):
                format_info.spacing = paragraph.spacing
            if hasattr(paragraph, 'margin_left'):
                format_info.margin_left = paragraph.margin_left
            if hasattr(paragraph, 'margin_right'):
                format_info.margin_right = paragraph.margin_right
        return format_info
    except Exception as e:
        logging.warning(f"格式信息获取失败: {e}")
        return pptTextFormat()

def apply_text_format(run, format_info, shape=None, text=None):
    """应用文本格式（增强颜色应用）"""
    try:
        if format_info.font_name:
            run.font.name = format_info.font_name
        if format_info.font_size:
            run.font.size = Pt(format_info.font_size)
        if format_info.font_bold is not None:
            run.font.bold = format_info.font_bold
        if format_info.font_italic is not None:
            run.font.italic = format_info.font_italic
        if format_info.font_underline is not None:
            run.font.underline = format_info.font_underline
        
        # 颜色处理逻辑增强
        if format_info.color:
            if isinstance(format_info.color, RGBColor):
                run.font.color.rgb = format_info.color
            elif isinstance(format_info.color, int):  # 处理主题颜色
                run.font.color.theme_color = format_info.color
            else:
                run.font.color.rgb = RGBColor(0, 0, 0)  # 默认黑色
        
        if hasattr(run, '_p'):
            paragraph = run._p
            if format_info.alignment is not None:
                paragraph.alignment = format_info.alignment
            if format_info.spacing is not None:
                paragraph.spacing = format_info.spacing
            if format_info.margin_left is not None:
                paragraph.margin_left = format_info.margin_left
            if format_info.margin_right is not None:
                paragraph.margin_right = format_info.margin_right
    except Exception as e:
        logging.warning(f"格式应用失败: {e}")

def translate_text_with_format(translation_core, text, source_lang, target_lang):
    """翻译文本（保持原逻辑）"""
    if not text.strip():
        return ""
    text_chunks = split_text(text)
    translated_chunks = []
    for chunk in text_chunks:
        translated_chunk = translation_core.translate_text(chunk, source_lang, target_lang)
        translated_chunks.append(translated_chunk)
    return "".join(translated_chunks)

def get_separators(language):
    """获取分隔符（保持原逻辑）"""
    if language == "Chinese":
        return '．，、：；！？”“‘’（）《》【】——…'
    elif language == "English":
        return ', : ; ! ? " \' ( ) - ...'
    elif language == "Japanese":
        return '．、：；！？”）（）『』【】ー…'
    else:
        return '. , : ;'

def split_text_into_parts(text, num_parts, target_lang):
    """分割文本（保持原逻辑）"""
    avg_length = len(text) // num_parts
    split_texts = []
    start = 0
    separators_all = get_separators(target_lang)
    separators = re.compile(r'[' + separators_all + ']')
    for i in range(num_parts):
        if i == num_parts - 1:
            split_texts.append(text[start:])
        else:
            end = start + avg_length
            match = separators.search(text, end)
            if match:
                end = match.end()
            else:
                end = start + avg_length
            split_texts.append(text[start:end])
            start = end
    return split_texts

# ----------------------- 辅助函数 -----------------------
def scan_shape(shape):
    """扫描形状（保持原逻辑）"""
    total = 0
    if hasattr(shape, "has_text_frame") and shape.has_text_frame:
        for paragraph in shape.text_frame.paragraphs:
            if paragraph.runs:
                for run in paragraph.runs:
                    if run.text.strip():
                        total += len(run.text)
            else:
                if paragraph.text.strip():
                    total += len(paragraph.text)
    if hasattr(shape, "has_table") and shape.has_table:
        for row in shape.table.rows:
            for cell in row.cells:
                if cell.text.strip():
                    total += len(cell.text)
    if hasattr(shape, "shapes"):
        for sub_shape in shape.shapes:
            total += scan_shape(sub_shape)
    return total

def process_text_frame(text_frame, container, translation_core, source_lang, target_lang, task, current_work, total_work):
    """处理文本框架（增强颜色保持）"""
    paragraph_formats = []
    for paragraph in text_frame.paragraphs:
        alignment = paragraph.alignment
        runs_info = []
        combined_text = ""
        format_infos = []
        for run in paragraph.runs:
            if run.text:
                combined_text += run.text
                # 获取每个run的完整格式信息
                format_infos.append(get_text_format(run))
        if combined_text.strip():
            translated_text = translate_text_with_format(translation_core, combined_text, source_lang, target_lang)
            # 保持原始run的数量和格式对应
            split_texts = split_text_into_parts(translated_text, len(format_infos), target_lang)
            runs_info.append((split_texts, format_infos))
            current_work[0] += len(combined_text)
            progress = (current_work[0] / total_work) * 100
            if task is not None:
                task.update_state(
                    state='PROGRESS',
                    meta={
                        'current': current_work[0],
                        'total': total_work,
                        'progress': round(progress, 1)
                    }
                )
        paragraph_formats.append((alignment, runs_info))
    
    # 清空文本但保留段落结构
    while len(text_frame.paragraphs) > 1:
        p = text_frame.paragraphs[-1]
        p._element.getparent().remove(p._element)
    first_paragraph = text_frame.paragraphs[0]
    first_paragraph.text = ""
    
    # 重新应用带格式的文本
    for idx, (alignment, runs_info) in enumerate(paragraph_formats):
        if idx == 0:
            paragraph = text_frame.paragraphs[0]
        else:
            paragraph = text_frame.add_paragraph()
        if alignment is not None:
            paragraph.alignment = alignment
            if runs_info and len(runs_info) > 0 and hasattr(paragraph, 'margin_left'):
                paragraph.margin_left = runs_info[0][1][0].margin_left
        for split_texts, format_infos in runs_info:
            for split_text, format_info in zip(split_texts, format_infos):
                run = paragraph.add_run()
                run.text = split_text
                # 应用原始run的格式（包含颜色）
                apply_text_format(run, format_info, container, split_text)

def process_table(table, translation_core, source_lang, target_lang, task, current_work, total_work):
    """处理表格（保持原逻辑）"""
    for row in table.rows:
        for cell in row.cells:
            if cell.text.strip():
                process_text_frame(cell.text_frame, cell, translation_core, source_lang, target_lang, task, current_work, total_work)

def process_shape(shape, translation_core, source_lang, target_lang, task, current_work, total_work):
    """处理形状（保持原逻辑）"""
    if hasattr(shape, "has_text_frame") and shape.has_text_frame:
        process_text_frame(shape.text_frame, shape, translation_core, source_lang, target_lang, task, current_work, total_work)
    elif hasattr(shape, "has_table") and shape.has_table:
        process_table(shape.table, translation_core, source_lang, target_lang, task, current_work, total_work)
    if hasattr(shape, "shapes"):
        for sub_shape in shape.shapes:
            process_shape(sub_shape, translation_core, source_lang, target_lang, task, current_work, total_work)

# ----------------------- 后处理 -----------------------
def adjust_text_frame_font_size(text_frame, container):
    """
    改进后的字体调整逻辑：
    1. 使用加权平均计算字体大小
    2. 更精确的字符宽度估算
    3. 强制最小字号限制
    """
    total_text = ""
    total_chars = 0
    sum_font_size = 0.0
    
    # 第一次遍历：收集文本信息和计算加权平均字号
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            run_text = run.text.strip()
            if run_text:
                # 计算每个run的字符数
                char_count = len(run_text)
                total_chars += char_count
                total_text += run_text
                
                # 获取当前run的字号（没有设置时使用默认12pt）
                run_size = 12.0  # 默认值
                if run.font.size and run.font.size.pt:
                    run_size = run.font.size.pt
                sum_font_size += run_size * char_count  # 加权计算

    # 计算加权平均字号
    if total_chars == 0:
        return
    avg_font_size = sum_font_size / total_chars

    # 获取容器尺寸（转换为pt）
    try:
        width_pt = container.width / 12700.0  # EMU转pt
        height_pt = container.height / 12700.0
    except AttributeError:
        return

    # 字符宽度估算（根据实际测试调整系数）
    avg_char_width = avg_font_size * 0.6  # 原1.2改为0.6更符合实际
    
    # 计算每行字符数（至少1个字符）
    chars_per_line = max(1, int(width_pt / avg_char_width))
    
    # 计算总行数（向上取整）
    total_lines = math.ceil(total_chars / chars_per_line)
    
    # 计算行高（包含行间距）
    line_height = avg_font_size * 1.2  # 1.2倍行距
    
    # 计算总需求高度
    required_height = total_lines * line_height
    
    # 计算缩放因子
    scaling_factor = 1.0
    if required_height > height_pt:
        scaling_factor = height_pt / required_height
    
    # 应用二次缩放（更保守的缩放）
    scaling_factor *= 0.95  # 预留5%的边距
    
    # 最小字号限制
    min_font_size = 8.0
    
    # 第二次遍历：应用缩放
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            if run.text.strip():
                # 获取原始字号或使用平均字号
                original_size = avg_font_size  # 默认使用平均字号
                if run.font.size and run.font.size.pt:
                    original_size = run.font.size.pt
                
                # 计算新字号
                new_size = original_size * scaling_factor
                new_size = max(new_size, min_font_size)
                
                # 应用新字号
                run.font.size = Pt(new_size)

def adjust_shape_font_size(shape):
    """形状字体调整（增加异常处理）"""
    try:
        if not hasattr(shape, 'width') or not hasattr(shape, 'height'):
            return
        
        # 文本框处理
        if hasattr(shape, "has_text_frame") and shape.has_text_frame:
            adjust_text_frame_font_size(shape.text_frame, shape)
        
        # 表格处理
        if hasattr(shape, "has_table") and shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    adjust_text_frame_font_size(cell.text_frame, cell)
        
        # 递归处理组合形状
        if hasattr(shape, "shapes"):
            for sub_shape in shape.shapes:
                adjust_shape_font_size(sub_shape)
    except Exception as e:
        logging.warning(f"字体调整失败: {str(e)}")                

def adjust_font_size_for_all_shapes(prs):
    """统一调整字体（保持原逻辑）"""
    for slide in prs.slides:
        for shape in slide.shapes:
            adjust_shape_font_size(shape)

# ----------------------- 主函数 -----------------------
def translate_powerpoint(translation_core, file_path, output_path, source_lang, target_lang, task):
    """主翻译函数"""
    logging.info(f"开始翻译PPT文件: {file_path}")
    try:
        prs = Presentation(file_path)
        total_work = 0
        for slide in prs.slides:
            for shape in slide.shapes:
                total_work += scan_shape(shape)
        current_work = [0]
        for slide in prs.slides:
            for shape in slide.shapes:
                process_shape(shape, translation_core, source_lang, target_lang, task, current_work, total_work)
        adjust_font_size_for_all_shapes(prs)
        prs.save(output_path)
        logging.info(f"PPT翻译完成: {output_path}")
    except Exception as e:
        logging.error(f"文件翻译失败 {file_path}: {e}")
        logging.error(traceback.format_exc())
        if task is not None:
            task.update_state(
                state='FAILURE',
                meta={
                    'exc_type': type(e).__name__,
                    'exc_message': str(e),
                    'traceback': traceback.format_exc()
                }
            )
        raise