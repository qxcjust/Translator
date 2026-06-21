"""文档解析器：支持 docx, txt, xlsx, pptx 的基础段落提取。
返回分段列表，每段为 dict: {id, type, text, meta}
"""
from pathlib import Path
from typing import List, Dict
import logging

from docx import Document
from pptx import Presentation
import openpyxl

logger = logging.getLogger(__name__)


def parse_document(path: str) -> List[Dict]:
    p = Path(path)
    suffix = p.suffix.lower()
    if suffix == ".docx":
        return _parse_docx(p)
    if suffix in (".ppt", ".pptx"):
        return _parse_pptx(p)
    if suffix in (".xls", ".xlsx"):
        return _parse_xlsx(p)
    if suffix == ".txt":
        return _parse_txt(p)
    # fallback: read as text
    return _parse_txt(p)


def _parse_docx(p: Path):
    doc = Document(p)
    segments = []
    for i, para in enumerate(doc.paragraphs):
        text = para.text.strip()
        if not text:
            continue
        seg = {"id": f"docx_{i}", "type": "paragraph", "text": text, "meta": {"source": str(p)}}
        segments.append(seg)
    logger.info(f"Parsed {len(segments)} paragraphs from {p}")
    return segments


def _parse_pptx(p: Path):
    prs = Presentation(p)
    segments = []
    idx = 0
    for sidx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = shape.text.strip()
                if text:
                    segments.append({"id": f"ppt_{sidx}_{idx}", "type": "slide_text", "text": text, "meta": {"slide": sidx}})
                    idx += 1
    logger.info(f"Parsed {len(segments)} text blocks from {p}")
    return segments


def _parse_xlsx(p: Path):
    wb = openpyxl.load_workbook(p, data_only=True)
    segments = []
    idx = 0
    for sheetname in wb.sheetnames:
        ws = wb[sheetname]
        for row in ws.iter_rows(values_only=True):
            for cell in row:
                if isinstance(cell, str) and cell.strip():
                    segments.append({"id": f"xls_{sheetname}_{idx}", "type": "cell", "text": cell.strip(), "meta": {"sheet": sheetname}})
                    idx += 1
    logger.info(f"Parsed {len(segments)} cells from {p}")
    return segments


def _parse_txt(p: Path):
    with p.open("r", encoding="utf-8", errors="ignore") as f:
        text = f.read()
    paragraphs = [seg.strip() for seg in text.splitlines() if seg.strip()]
    return [{"id": f"txt_{i}", "type": "paragraph", "text": paragraphs[i], "meta": {"source": str(p)}} for i in range(len(paragraphs))]
