import os
from docx import Document
from docx.table import Table
from docx.text.paragraph import Paragraph
from docx.shared import Pt
from pptx.dml.color import RGBColor

def update_progress(task, current_work, total_work):
    """更新翻译进度"""
    if task is not None:
        progress = (current_work / total_work) * 100 if total_work > 0 else 0
        task.update_state(
            state='PROGRESS',
            meta={
                'current': current_work,
                'total': total_work,
                'progress': round(progress, 1)
            }
        )

def translate_paragraphs(translation_core, paragraphs, source_lang, target_lang, task, current_work, total_work):
    """
    逐段落、逐 Run 翻译，并更新 current_work。
    :param paragraphs: doc.paragraphs 或 cell.paragraphs 或 shape.text_frame.paragraphs
    :return: 更新后的 current_work
    """
    for para in paragraphs:
        for run in para.runs:
            text = run.text
            if text.strip():
                translated_text = translation_core.translate_text(text, source_lang, target_lang)
                run.text = translated_text
                current_work += len(text)
                update_progress(task, current_work, total_work)
    return current_work

def extract_table_text(translation_core, table, source_lang, target_lang, task, current_work, total_work, level=0):
    """
    递归提取表格中的文本，包括嵌套表格。
    采用逐段落、逐 Run 的方式翻译，以保留格式并确保不遗漏。
    :param table: Word 文档中的表格对象
    :param level: 嵌套层级（用于调试或打印）
    :return: 更新后的 current_work
    """
    for row in table.rows:
        for cell in row.cells:
            # 翻译单元格中的所有段落
            current_work = translate_paragraphs(
                translation_core, 
                cell.paragraphs, 
                source_lang, 
                target_lang, 
                task, 
                current_work, 
                total_work
            )
            # 添加调试信息
            print(f"翻译单元格内容: {cell.text}，当前进度: {current_work}/{total_work}")
            # 检查单元格是否包含嵌套表格
            if cell.tables:
                for nested_table in cell.tables:
                    current_work = extract_table_text(
                        translation_core, 
                        nested_table, 
                        source_lang, 
                        target_lang, 
                        task, 
                        current_work, 
                        total_work, 
                        level + 1
                    )
    return current_work

def calculate_total_work(doc):
    """计算文档总工作量（包括页眉、页脚、正文）"""
    total_work = 0

    # 处理正文内容
    for para in doc.paragraphs:
        total_work += len(para.text)
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                total_work += len(cell.text)
                # 检查单元格是否包含嵌套表格
                if cell.tables:
                    for nested_table in cell.tables:
                        for nested_row in nested_table.rows:
                            for nested_cell in nested_row.cells:
                                total_work += len(nested_cell.text)

    # 处理内联形状的文本
    for shape in doc.inline_shapes:
        if hasattr(shape, 'text_frame'):
            text_frame = shape.text_frame
            for para in text_frame.paragraphs:
                total_work += len(para.text)

    # 处理页眉和页脚
    for section in doc.sections:
        # 页眉
        header = section.header
        for para in header.paragraphs:
            total_work += len(para.text)
        for table in header.tables:
            for row in table.rows:
                for cell in row.cells:
                    total_work += len(cell.text)
                    # 检查单元格是否包含嵌套表格
                    if cell.tables:
                        for nested_table in cell.tables:
                            for nested_row in nested_table.rows:
                                for nested_cell in nested_row.cells:
                                    total_work += len(nested_cell.text)
        # 页脚
        footer = section.footer
        for para in footer.paragraphs:
            total_work += len(para.text)
        for table in footer.tables:
            for row in table.rows:
                for cell in row.cells:
                    total_work += len(cell.text)
                    # 检查单元格是否包含嵌套表格
                    if cell.tables:
                        for nested_table in cell.tables:
                            for nested_row in nested_table.rows:
                                for nested_cell in nested_row.cells:
                                    total_work += len(nested_cell.text)
    return total_work

def translate_word(translation_core, file_path, output_path, source_lang, target_lang, task):
    """
    翻译 Word 文档的全部内容（正文、页眉、页脚）
    """
    doc = Document(file_path)
    
    # 计算总文本量（包括页眉页脚）
    total_work = calculate_total_work(doc)
    current_work = 0

    # 1. 翻译正文内容
    # 1.1 翻译段落
    current_work = translate_paragraphs(
        translation_core, 
        doc.paragraphs, 
        source_lang, 
        target_lang, 
        task, 
        current_work, 
        total_work
    )

    # 1.2 翻译正文中的表格
    for table in doc.tables:
        current_work = extract_table_text(
            translation_core, 
            table, 
            source_lang, 
            target_lang, 
            task, 
            current_work, 
            total_work
        )

    # 1.3 翻译内联形状
    for shape in doc.inline_shapes:
        if not hasattr(shape, 'text_frame'):
            continue
        text_frame = shape.text_frame
        current_work = translate_paragraphs(
            translation_core, 
            text_frame.paragraphs, 
            source_lang, 
            target_lang, 
            task, 
            current_work, 
            total_work
        )

    # 2. 翻译页眉页脚
    for section in doc.sections:
        # 2.1 处理页眉
        header = section.header
        # 翻译页眉段落
        current_work = translate_paragraphs(
            translation_core, 
            header.paragraphs, 
            source_lang, 
            target_lang, 
            task, 
            current_work, 
            total_work
        )
        # 翻译页眉表格
        for table in header.tables:
            current_work = extract_table_text(
                translation_core, 
                table, 
                source_lang, 
                target_lang, 
                task, 
                current_work, 
                total_work
            )

        # 2.2 处理页脚
        footer = section.footer
        # 翻译页脚段落
        current_work = translate_paragraphs(
            translation_core, 
            footer.paragraphs, 
            source_lang, 
            target_lang, 
            task, 
            current_work, 
            total_work
        )
        # 翻译页脚表格
        for table in footer.tables:
            current_work = extract_table_text(
                translation_core, 
                table, 
                source_lang, 
                target_lang, 
                task, 
                current_work, 
                total_work
            )

    # 保存文档
    doc.save(output_path)